#import torch
#from transformers import BertForQuestionAnswering
#from transformers import BertTokenizer
import re
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import fitz
from docx import Document
from pptx import Presentation
import pandas as pd
from fuzzywuzzy import fuzz
import os

class qnatb(object):
    stopwords = ['i', 'me', 'my', 'myself', 'we', 'our', 'ours', 'ourselves', 'you', 'your', 'yours', 'yourself',
                 'yourselves', 'he', 'him', 'his', 'himself', 'she', 'her', 'hers', 'herself', 'it', 'its', 'itself',
                 'they', 'them', 'their', 'theirs', 'themselves', 'what', 'which', 'who', 'whom', 'this', 'that',
                 'these', 'those', 'am', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had',
                 'having', 'do', 'does', 'did', 'doing', 'a', 'an', 'the', 'and', 'but', 'if', 'or', 'because', 'as',
                 'until', 'while', 'of', 'at', 'by', 'for', 'with', 'about', 'against', 'between', 'into', 'through',
                 'during', 'before', 'after', 'above', 'below', 'to', 'from', 'up', 'down', 'in', 'out', 'on', 'off',
                 'over', 'under', 'again', 'further', 'then', 'once', 'here', 'there', 'when', 'where', 'why', 'how',
                 'all', 'any', 'both', 'each', 'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not',
                 'only', 'own', 'same', 'so', 'than', 'too', 'very', 's', 't', 'can', 'will', 'just', 'don', 'should',
                 'now']

    def __init__(self):
        pass
        #self.model_path = model_path
        #self.model = BertForQuestionAnswering.from_pretrained(model_path)
        #self.tokenizer = BertTokenizer.from_pretrained(model_path)


    @staticmethod
    def clean(sent):
        sent = re.sub(r'<.*?>', " ", sent)  # html tags
        sent = re.sub(r'h\s*t\s*t\s*p\s*s?://\S+|www\.\S+', " ", sent)  # remove urls

        sent = re.sub('\n', " ", sent)
        sent = re.sub(r'\(.*?\)', " ", sent)  # all inside ()
        sent = re.sub(r'\[.*?\]', " ", sent)  # all inside ()
        sent = re.sub(r'[^A-Za-z0-9\.,\?\(\)\[\]\/ ]', " ", sent)
        sent = re.sub('\s+', " ", sent)
        return sent

    @staticmethod
    def split_into_sentences(text):
        alphabets = "([A-Za-z])"
        prefixes = "(Mr|St|Mrs|Ms|Dr|No)[.]"
        suffixes = "(Inc|Ltd|Jr|Sr|Co)"
        starters = "(Mr|Mrs|Ms|Dr|He\s|She\s|It\s|They\s|Their\s|Our\s|We\s|But\s|However\s|That\s|This\s|Wherever)"
        acronyms = "([A-Z][.][A-Z][.](?:[A-Z][.])?)"
        decimals = "(\d*[.]\d*)"
        websites = "[.](com|net|org|io|gov|co|in)"
        decimals = r'\d\.\d'

        text = " " + text + "  "
        text = text.replace("\n", " ")
        text = re.sub(decimals, lambda g: re.sub(r'\.', '<prd>', g[0]), text)
        # text= re.sub(r'(?<=\[).+?(?=\])', "", text) # remove everything inside square brackets
        # text= re.sub(r'(?<=\().+?(?=\))', "", text) # remove everything inside  brackets
        text = re.sub(r'\[(\w*)\]', "", text)  # remove evrything in sq brackets with sq brackets
        text = re.sub(prefixes, "\\1<prd>", text)
        text = re.sub(websites, "<prd>\\1", text)
        if "i.e." in text: text = text.replace("i.e.", "i<prd>e<prd>")
        if "e.g" in text: text = text.replace("e.g", "e<prd>g")
        if "www." in text: text = text.replace("www.", "www<prd>")
        text = re.sub("\s" + alphabets + "[.] ", " \\1<prd> ", text)
        text = re.sub(acronyms + " " + starters, "\\1<stop> \\2", text)
        text = re.sub(alphabets + "[.]" + alphabets + "[.]" + alphabets + "[.]", "\\1<prd>\\2<prd>\\3<prd>", text)
        text = re.sub(alphabets + "[.]" + alphabets + "[.]", "\\1<prd>\\2<prd>", text)
        text = re.sub(" " + suffixes + "[.] " + starters, " \\1<stop> \\2", text)
        text = re.sub(" " + suffixes + "[.]", " \\1<prd>", text)
        text = re.sub(" " + alphabets + "[.]", " \\1<prd>", text)
        if "”" in text: text = text.replace(".”", "”.")
        if "\"" in text: text = text.replace(".\"", "\".")
        if "!" in text: text = text.replace("!\"", "\"!")
        if "?" in text: text = text.replace("?\"", "\"?")
        text = text.replace(".", ".<stop>")
        text = text.replace("?", "?<stop>")
        text = text.replace("!", "!<stop>")
        text = text.replace("<prd>", ".")
        sentences = text.split("<stop>")
        # sentences = sentences[:-1]

        sentences = [s.strip() for s in sentences]
        return sentences
    
    @staticmethod
    def tb_index_pdf(Folder, file, tb_index):
        
        doc = fitz.open(os.path.join(Folder, file))
        try:
            metadata= {i:j for i,j in doc.metadata.items() if i in ["format", "title", "author", "creationDate", "modDate"]}
        except:
            metadata={'format':'PDF', 'title': "", "author": "", "creationDate":"", "modDate":""}
            
        metadata['filename']=file

        for num, page in enumerate(doc):
            try:
                text = page.getText().encode('utf8')
                text = text.decode('utf8')
                text = qnatb.clean(text)
                sentences = qnatb.split_into_sentences(text)
                for sent in sentences:
                    tb_index.append({
                        "doc": file.split('\\')[-1],
                        "page": num,
                        "sentence": sent
    
                    })
            except:
                tb_index.append({
                    "doc": file.split('\\')[-1],
                    "page": num,
                    "sentence": ""
    
                })
        return tb_index, metadata

    @staticmethod
    def tb_index_docx(Folder, file, tb_index):
        doc= Document(os.path.join(Folder, file))
        try:
            metadata={}
            prop=doc.core_properties
            metadata['format']="docx"
            metadata['title']=prop.subject
            metadata['author']=prop.author
            metadata['creationDate']=prop.created
            metadata['modDate']=prop.modified
        except:
            metadata={'format':'docx', 'title': "", "author": "", "creationDate":"", "modDate":""}
                        
        metadata['filename']=file

        text=[]
        try:
            for para in doc.paragraphs:
                text.append(para.text)
            text=" ".join([i for i in text if i.strip()!=""])
            #text=text.encode('utf8')
            text = qnatb.clean(text)
            sentences = qnatb.split_into_sentences(text)
            for sent in sentences:
                tb_index.append({
                    "doc": file.split('\\')[-1],
                    "page": "-",
                    "sentence": sent
        
                })
        except:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": "-",
                "sentence": "Not read"
    
            })
        return tb_index, metadata

    @staticmethod
    def tb_index_pptx(Folder, file, tb_index):
        ppt= Presentation(os.path.join(Folder, file))
        
        metadata={'format':'pptx', 'title': "", "author": "", "creationDate":"", "modDate":""}
        metadata['filename']=file
        
        for num, slide in enumerate(ppt.slides):
            try:
                all_text=[]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text= shape.text
                        text = qnatb.clean(text)
                        sentences = qnatb.split_into_sentences(text)
                        [all_text.append(i) for i in sentences]
                all_text= " ".join(i for i in all_text)        
                tb_index.append({
                    "doc": file.split('\\')[-1],
                    "page": num,
                    "sentence": all_text
                })
    
            except:
                tb_index.append({
                    "doc": file.split('\\')[-1],
                    "page": num,
                    "sentence": ""
    
                })
        return tb_index, metadata


    def files_processor_tb(self, Folder):
        tb_index = []
        
        response_file_processing=[]
        
        metadata_all=[]
        
        files= [i for i in os.listdir(Folder)]
        
        
        for file in files:
            
            if file.endswith('pdf'):
                try:
                    _, metadata= qnatb.tb_index_pdf(Folder, file, tb_index)
                    metadata_all.append(metadata)
                    response= {"filename": file, "msg":"", "status": "success"}
                    response_file_processing.append(response)
                except Exception as e:
                    response= {"filename": file, "msg":str(e), "status": "failed"}
                    response_file_processing.append(response)
                    
                    
            elif file.endswith('docx'):
                try:
                    _, metadata= qnatb.tb_index_docx(Folder, file, tb_index)
                    metadata_all.append(metadata)
                    response= {"filename": file, "msg":"", "status": "success"}
                    response_file_processing.append(response)
                except Exception as e:
                    response= {"filename": file, "msg":str(e), "status": "failed"}
                    response_file_processing.append(response)
                
            elif file.endswith('pptx'):
                try:
                    _, metadata= qnatb.tb_index_pptx(Folder, file, tb_index)
                    metadata_all.append(metadata)
                    response= {"filename": file, "msg":"", "status": "success"}
                    response_file_processing.append(response)
                except Exception as e:
                    response= {"filename": file, "msg":str(e), "status": "failed"}
                    response_file_processing.append(response)
            else:
                print(file)

            
        self.metadata_all= metadata_all
        self.tb_index=tb_index
        all_sents= [i['sentence'] for i in tb_index]
        self.all_sents= all_sents
        
        
        #run metadata
        self.metadata()
        return tb_index, all_sents, response_file_processing
    
    
    def reg_ind(self, words):
        if "," in words:
            #words= words.split(',')
            words= [i.strip().lower() for i in words.split(",")]
            reg= "|".join(words)
            tb_index_reg=self.tb_index
            tb_index_reg=[i for i in self.tb_index if len(re.findall(reg, i['sentence'].lower()))>0]
            
        elif "+" in words:
            words= [i.strip().lower() for i in words.split("+")]
            tb_index_reg=self.tb_index
            for word in words:   
                tb_index_reg=[i for i in tb_index_reg if len(re.findall(word, i['sentence'].lower()))>0]
        else:
            words= words.strip().lower()
            tb_index_reg=[i for i in self.tb_index if len(re.findall(words, i['sentence'].lower()))>0]
        
        
        docs= list(set([i['doc'] for i in tb_index_reg]))
        
        overall_dict={i:sum([1 for j in tb_index_reg if j['doc']==i]) for i in docs}
        #number of sentences not occurances
        
        return tb_index_reg, overall_dict, docs
    
    @staticmethod
    def extract_doc_reg_index(tb_index_reg, doc):
        reg_tb_index= [i for i in tb_index_reg if i['doc']==doc]
        req_df= pd.DataFrame(reg_tb_index)
        req_df.drop('doc', axis=1, inplace=True)
        return req_df
    
    @staticmethod
    def get_score(question_tfidf, sent):
        scr=0
        counter=0
        sent = " ".join([i for i in sent.lower().split() if i not in qnatb.stopwords])
        
        for token in question_tfidf.split():
            scr+=fuzz.partial_ratio(token,sent)
            counter+=1
        return scr/counter          
            
            
    def get_response_sents(self, question, max_length=None):

        question = " ".join([i for i in question.split() if i not in qnatb.stopwords])
        
        dict_scores={num:qnatb.get_score(question,i['sentence']) for num,i in enumerate(self.tb_index)}
        
        dict_scores = {k: v for k, v in sorted(dict_scores.items(), key=lambda item: item[1], reverse=True)}
        
        final_response_dict = [self.tb_index[i] for i, j in dict_scores.items()]

        if max_length:
            final_response_dict = [i for i in final_response_dict if len(i['sentence'].split()) > max_length]

        return final_response_dict

    # def answer_question(self, question, answer_text):
    #     question = qnatb.clean(question)
    #     encoded_dict = self.tokenizer.encode_plus(text=question, text_pair=answer_text, add_special=True)
    #     input_ids = encoded_dict['input_ids']
    #     segment_ids = encoded_dict['token_type_ids']
    #     assert len(segment_ids) == len(input_ids)
    #     output = self.model(torch.tensor([input_ids]),  # The tokens representing our input text.
    #                         token_type_ids=torch.tensor(
    #                             [segment_ids]))  # The segment IDs to differentiate question from answer_text

    #     answer_start = torch.argmax(output['start_logits'])
    #     start_logit = output['start_logits'][0][answer_start].detach().numpy()
    #     answer_end = torch.argmax(output['end_logits'])

    #     tokens = self.tokenizer.convert_ids_to_tokens(input_ids)
    #     answer = tokens[answer_start]
    #     for i in range(answer_start + 1, answer_end + 1):

    #         if tokens[i][0:2] == '##':
    #             answer += tokens[i][2:]
    #         else:
    #             answer += ' ' + tokens[i]
    #     return answer, start_logit

    # def retrieve_answer(self, question, top=10, max_length=None):

    #     response_sents = self.get_response_sents(question, max_length=max_length)
    #     max_logit = 3
    #     logits = []
    #     correct_answer = "Please rephrase"
    #     answer_extracted = "Please rephrase"

    #     for num, answer_text in enumerate(response_sents[0:top]):
    #         answer, start_logit = self.answer_question(question, answer_text['sentence'])
    #         logits.append(start_logit)
    #         if start_logit > max_logit:
    #             max_logit = start_logit
    #             correct_answer = answer
    #             answer_extracted = answer_text
    #             # answer_num=num

    #     return correct_answer, answer_extracted, max_logit, logits

    def get_top_n(self, question, top=10, max_length=None):

        response_sents = self.get_response_sents(question, max_length=max_length)
        # top_responses = []

        # for num, answer_text in enumerate(response_sents[0:top]):
        #     answer, start_logit = self.answer_question(question, answer_text['sentence'])
        #     top_response = response_sents[num]
        #     top_response['start_logit'] = start_logit
        #     top_response['answer'] = answer
        #     top_responses.append(top_response)
        # top_responses = sorted(top_responses, key=lambda item: item['start_logit'], reverse=True)
        # responses = top_responses + response_sents[top:]
        # return responses
        return response_sents
    
    def stats(self):
        docs= list(set([i['doc'] for i in self.tb_index]))
        stats=[]
        for doc in docs:
            st={}
            st['doc']=doc
            st['pages']=len(set([i['page'] for i in self.tb_index if i['doc']==doc]))
            st['words']= sum([len(i['sentence'].split()) for i in self.tb_index if i['doc']==doc])
            stats.append(st)
        return stats
    
    def metadata(self):
        try:
            md= self.metyadata_all
            updated_md=[]
            for metadata_file in md:
                try:
                    filename= metadata_file['filename'].split('/')[-1]
                    metadata_file['pages']=len(set([i['page'] for i in self.tb_index if i['filename']==filename]))
                    metadata_file['words']=sum([len(i['sentence'].split()) for i in self.tb_index if i['filename']==filename])
                except:
                    metadata_file['pages']=None
                    metadata_file['words']=None

                updated_md.append(metadata_file)
            self.metadata_all= updated_md
        except:
            pass
        
    
            
        



# qna = qnatb(model_path=r'C:\Users\ELECTROBOT\Desktop\Bert-qa\model')
# import os
# files= os.listdir(r"C:\Users\ELECTROBOT\Desktop\Desktop_kachra")

# names=[os.path.join(r"C:\Users\ELECTROBOT\Desktop\Desktop_kachra",i) for i in files]
# _, _, _, _ = qna.files_processor_tb(names)

# search_data="who did federer marry"
# responses = qna.get_top_n(search_data, top=10, max_length=7)










