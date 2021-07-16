# -*- coding: utf-8 -*-
"""
Created on Wed Jul 14 11:26:18 2021

@author: ELECTROBOT
"""
import re
import fitz
from docx import Document
from pptx import Presentation


def clean(text, html=False,brackets=False, digits=False, only_english=False, lower=False):
    
    #for xml keep this true
    if html==True:
        text=re.sub(r'<.*?>|&([a-z0-9]+|#[0-9]{1,6}|#x[0-9a-f]{1,6});'," ", text)
    #remove everything in brackets
    if brackets== True:
        text= re.sub(r'(?<=\[).+?(?=\])', "", text)
        text= re.sub(r'(?<=\().+?(?=\))', "", text)
    #remove digits
    #if digits== True:
    #    text= re.sub(r'\d+\.\d+|\d+', " ", text)

    if only_english==True:
        text=re.sub(r'[^A-Za-z ]', " ", text)
    
    if lower==True:
        text=text.lower()

    #text= re.sub(r'[^A-Za-z0-9\.\(\)\[\]@ ]', " ", text)
    text= re.sub(r'\s+', " ", text.strip())
    
    return text
    
    
def split_into_sentences(text):
    
    alphabets = "([A-Za-z])"
    prefixes = "(Mr|St|Mrs|Ms|Dr|No|no)[.]"
    suffixes = "(Inc|Ltd|Jr|Sr|Co)"
    starters = "(Mr|Mrs|Ms|Dr|He\s|She\s|It\s|They\s|Their\s|Our\s|We\s|But\s|However\s|That\s|This\s|Wherever)"
    acronyms = "([A-Z][.][A-Z][.](?:[A-Z][.])?)"
    decimals = "(\d*[.]\d*)"
    websites = "[.](com|net|org|io|gov|co|in)"
    decimals = r'\d\.\d'

    
    text = " " + text + "  "
    text = text.replace("\n", " ")
    text = re.sub(decimals, lambda g: re.sub(r'\.', '<prd>', g[0]), text)
    # text= re.sub(r'(?<=\[).+?(?=\])', "", text) # remove everything inside square brackets
    # text= re.sub(r'(?<=\().+?(?=\))', "", text) # remove everything inside  brackets
    text = re.sub(r'\[(\w*)\]', "", text)  # remove evrything in sq brackets with sq brackets
    text = re.sub(prefixes, "\\1<prd>", text)
    text = re.sub(websites, "<prd>\\1", text)
    if "i.e." in text: text = text.replace("i.e.", "i<prd>e<prd>")
    if "e.g." in text: text = text.replace("e.g.", "e<prd>g<prd>")
    if "e.g" in text: text = text.replace("e.g", "e<prd>g")
    if "www." in text: text = text.replace("www.", "www<prd>")
    text = re.sub("\s" + alphabets + "[.] ", " \\1<prd> ", text)
    text = re.sub(acronyms + " " + starters, "\\1<stop> \\2", text)
    text = re.sub(alphabets + "[.]" + alphabets + "[.]" + alphabets + "[.]", "\\1<prd>\\2<prd>\\3<prd>", text)
    text = re.sub(alphabets + "[.]" + alphabets + "[.]", "\\1<prd>\\2<prd>", text)
    text = re.sub(" " + suffixes + "[.] " + starters, " \\1<stop> \\2", text)
    text = re.sub(" " + suffixes + "[.]", " \\1<prd>", text)
    text = re.sub(" " + alphabets + "[.]", " \\1<prd>", text)
    if "”" in text: text = text.replace(".”", "”.")
    if "\"" in text: text = text.replace(".\"", "\".")
    if "!" in text: text = text.replace("!\"", "\"!")
    if "?" in text: text = text.replace("?\"", "\"?")
    
    text= re.sub("\w[.]\w", "<prd>", text) #replace dot in \w.\w by prd
    
    text = text.replace(".", ".<stop>")
    text = text.replace("?", "?<stop>")
    text = text.replace("!", "!<stop>")
    text = text.replace("<prd>", ".")
    sentences = text.split("<stop>")
    # sentences = sentences[:-1]

    sentences = [s.strip() for s in sentences]
    
    return sentences




def tb_index_pdf(file):
    tb_index=[]
    doc = fitz.open(file)
    for num, page in enumerate(doc):
        try:
            text = page.getText().encode('utf8')
            text = text.decode('utf8')
            text = clean(text)
            sentences = split_into_sentences(text)
            for sent in sentences:
                tb_index.append({
                    "doc": file.split('\\')[-1],
                    "page": num,
                    "sentence": sent

                })
        except:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": ""

            })
    return tb_index
    
def tb_index_docx(file):
    doc= Document(file)
    tb_index=[]
    text=[]
    try:
        for para in doc.paragraphs:
            text.append(para.text)
        text=" ".join([i for i in text if i.strip()!=""])
        #text=text.encode('utf8')
        text = clean(text)
        sentences = split_into_sentences(text)
        for sent in sentences:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": "-",
                "sentence": sent
    
            })
    except:
        tb_index.append({
            "doc": file.split('\\')[-1],
            "page": "-",
            "sentence": "Not read"

        })
    return tb_index



def tb_index_pptx(file):
    ppt= Presentation(file)
    tb_index=[]
    
    for num, slide in enumerate(ppt.slides):
        try:
            all_text=[]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text= shape.text
                    text = clean(text)
                    sentences = split_into_sentences(text)
                    [all_text.append(i) for i in sentences]
            all_text= " ".join(i for i in all_text)        
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": all_text
            })

        except:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": ""

            })
    return tb_index

                
                
                
                
    


def files_processor_tb(file):
        
    if file.endswith(".pdf"):
        try:
            tb_index=tb_index_pdf(file)
        except:
            print("error reading",file)
    elif file.endswith('.docx'):
        try:
            tb_index=tb_index_docx(file)
        except:
            print("error reading",file)
    elif file.endswith('.pptx'):
        try:
            tb_index=tb_index_pptx(file)
        except:
            print("error reading",file)
    else:
        print("error reading",file)
   

    return tb_index









stopwords = ['i', 'me', 'my', 'myself', 'we', 'our', 'ours', 'ourselves', 'you', 'your', 'yours', 'yourself',
                 'yourselves', 'he', 'him', 'his', 'himself', 'she', 'her', 'hers', 'herself', 'it', 'its', 'itself',
                 'they', 'them', 'their', 'theirs', 'themselves', 'what', 'which', 'who', 'whom', 'this', 'that',
                 'these', 'those', 'am', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had',
                 'having', 'do', 'does', 'did', 'doing', 'a', 'an', 'the', 'and', 'but', 'if', 'or', 'because', 'as',
                 'until', 'while', 'of', 'at', 'by', 'for', 'with', 'about', 'against', 'between', 'into', 'through',
                 'during', 'before', 'after', 'above', 'below', 'to', 'from', 'up', 'down', 'in', 'out', 'on', 'off',
                 'over', 'under', 'again', 'further', 'then', 'once', 'here', 'there', 'when', 'where', 'why', 'how',
                 'all', 'any', 'both', 'each', 'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not',
                 'only', 'own', 'same', 'so', 'than', 'too', 'very', 's', 't', 'can', 'will', 'just', 'don', 'should',
                 'now']

#def read(file):
    









    