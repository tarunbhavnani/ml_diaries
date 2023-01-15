# -*- coding: utf-8 -*-
"""
Created on Thu Jun 24 13:11:16 2021

@author: ELECTROBOT
"""
import re
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import fitz

from docx import Document

from pptx import Presentation


def ngrams(string):
    n=4
    string = re.sub(r'[,-./]|\sBD', r'', string)
    words=string.split()
    ngram=[]
    for word in words:
        #print(word)
        ngrams = zip(*[word[i:] for i in range(n)])
        ret=[''.join(ngram) for ngram in ngrams]
        [ngram.append(i) for i in ret]
    #[ngram.append(i) for i in string.split()]
    return ngram





import glob
files=glob.glob(r"C:\Users\ELECTROBOT\Desktop\da_test_files\*")


file=r"C:\Users\ELECTROBOT\Desktop\da_test_files\rf.pdf"


def tb_index_pdf(file, tb_index):
    doc = fitz.open(file)
    for num, page in enumerate(doc):
        try:
            text = page.getText().encode('utf8')
            text = text.decode('utf8')
            text = clean(text)
            sentences = split_into_sentences(text)
            for sent in sentences:
                tb_index.append({
                    "doc": file.split('\\')[-1],
                    "page": num,
                    "sentence": sent

                })
        except:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": ""

            })
    return tb_index
    
def tb_index_docx(file, tb_index):
    doc= Document(file)
    text=[]
    try:
        for para in doc.paragraphs:
            text.append(para.text)
        text=" ".join([i for i in text if i.strip()!=""])
        #text=text.encode('utf8')
        text = clean(text)
        sentences = split_into_sentences(text)
        for sent in sentences:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": "-",
                "sentence": sent
    
            })
    except:
        tb_index.append({
            "doc": file.split('\\')[-1],
            "page": "-",
            "sentence": "Not read"

        })
    return tb_index



def tb_index_pptx(file, tb_index):
    ppt= Presentation(path)
    
    for num, slide in enumerate(ppt.slides):
        try:
            all_text=[]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text= shape.text
                    text = clean(text)
                    sentences = split_into_sentences(text)
                    [all_text.append(i) for i in sentences]
            all_text= " ".join(i for i in all_text)        
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": all_text
            })

        except:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": ""

            })
    return tb_index

                
                
                
                
    


def files_processor_tb(files):
    tb_index = []
    for file in files:
        print(file)
        if file.endswith(".pdf"):
            try:
                tb_index_pdf(file, tb_index)
            except:
                print(file)
        elif file.endswith('.docx'):
            try:
                tb_index_docx(file, tb_index)
            except:
                print(file)
        elif file.endswith('.pptx'):
            try:
                tb_index_pptx(file, tb_index)
            except:
                print(file)
        else:
            print(file)

    
    all_sents = [i['sentence'] for i in tb_index]
    all_sents = all_sents
    vec = TfidfVectorizer(analyzer=ngrams)  # this performs much better but exact words

    vec.fit([i.lower() for i in all_sents])
    
    tfidf_matrix = vec.transform([i.lower() for i in all_sents])
    

    return tb_index, all_sents, vec, tfidf_matrix


tb_index, all_sents, vec, tfidf_matrix=files_processor_tb(files)




question= "who did federer married"
question_tfidf = " ".join([i for i in question.split() if i not in stopwords])
question_vec = vec.transform([question_tfidf])



scores = cosine_similarity(tfidf_matrix, question_vec)
scores = [i[0] for i in scores]
dict_scores = {i: j for i, j in enumerate(scores)}
dict_scores = {k: v for k, v in sorted(dict_scores.items(), key=lambda item: item[1], reverse=True)}
# get top n sentences
# final_response_dict=[self.tb_index[i] for i in dict_scores]
final_response_dict = [tb_index[i] for i, j in dict_scores.items() if j > 0.1]
final_response_dict = [i for i in final_response_dict if len(i['sentence'].split()) > 7]





# =============================================================================
# regex from tb_index
# =============================================================================

# reg= r'roger|married'

# reg=r'^(?=.*\bfederer\b)(?=.*\bmarried\b).*$'


# def reg_ind(word, tb_index):
#     #pass

#     tb_index_reg=[len(re.findall(word,i['sentence'].lower())) for i in tb_index]
#     #tb_index[0]['reg']=2

#     for num,i in enumerate(tb_index_reg):
#         tb_index[num]['reg']=i

#     tb_index_reg= [i for i in tb_index if i['reg']>0]
    
#     docs= list(set([i['doc'] for i in tb_index_reg]))
#     overall_dict={i:sum([j['reg'] for j in tb_index_reg if j['doc']==i]) for i in docs}
    
    
    
#     return tb_index_reg, overall_dict


# jkl,jkl1= reg_ind("flask", tb_index)




words= "federer , final, atp" #if any
words="federer + final + atp"# plus if all
words= "married "# nothing then search exact

def reg_ind(words, tb_index):
    if "," in words:
        #words= words.split(',')
        words= [i.strip() for i in words.split(",")]
        reg= "|".join(words)
        tb_index_reg=tb_index
        tb_index_reg=[i for i in tb_index if len(re.findall(reg, i['sentence']))>0]
        
    elif "+" in words:
        words= [i.strip() for i in words.split("+")]
        tb_index_reg=tb_index
        for word in words:   
            tb_index_reg=[i for i in tb_index_reg if len(re.findall(word, i['sentence']))>0]
    else:
        words= words.strip()
        tb_index_reg=[i for i in tb_index if len(re.findall(words, i['sentence']))>0]
    
    
    docs= list(set([i['doc'] for i in tb_index_reg]))
    
    overall_dict={i:sum([1 for j in tb_index_reg if j['doc']==i]) for i in docs}
    #number of sentences not occurances
    
    return tb_index_reg, overall_dict, docs
    


def extract_doc_reg_index(tb_index_reg, doc):
    req_tb_index=[i for i in tb_index_reg if i['doc']==doc]
    req_df=pd.DataFrame(req_tb_index)
    req_df.drop('doc', axis=1, inplace=True)
    #convert to html on put on webpage
    return req_df
    



tb_index_reg, overall_dict, docs=reg_ind(words, tb_index)

import pandas as pd
rf_ind=extract_doc_reg_index(tb_index_reg, 'rf.pdf')

        
        
    