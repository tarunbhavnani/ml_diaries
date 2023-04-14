# -*- coding: utf-8 -*-
"""
Created on Fri Jun 18 13:42:00 2021

@author: tarun
"""
import zipfile
import xml.etree.ElementTree

path=r"C:\Users\tarun\Desktop\test.docx"

WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
PARA = WORD_NAMESPACE + 'p'
TEXT = WORD_NAMESPACE + 't'
TABLE = WORD_NAMESPACE + 'tbl'
ROW = WORD_NAMESPACE + 'tr'
CELL = WORD_NAMESPACE + 'tc'

with zipfile.ZipFile(path) as docx:
    tree = xml.etree.ElementTree.XML(docx.read('word/document.xml'))

for table in tree.iter(TABLE):
    #pass
    for row in table.iter(ROW):
        for cell in row.iter(CELL):
            print (''.join(node.text for node in cell.iter(TEXT)))


import docx

doc= docx.Document(path)
result = [p.text for p in doc.paragraphs]


import docx2txt
 
# read in word file
result = docx2txt.process(path)
print(result)
result = docx2txt.process(path, r"C:\Users\tarun\Desktop\doc2txt")

docx2txt.properties


from docx2python import docx2python
 
# extract docx content
doc_result = docx2python(path)


# get separate components of the document
doc_result.body
 
# get the text from Zen of Python
doc_result.body[0]
 
# get the image
doc_result[1] 
 
# get the table text
doc_result[2]


import pandas as pd
 
 
pd.DataFrame(doc_result.body[1][1:]).applymap(lambda val: val[0].strip("\t"))


type(doc_result.images) # dict
 
for key,val in doc_result.images.items():
    f = open(key, "wb")
    f.write(val)
    f.close()

#all data
text=doc_result.text
print(text)
import re

def cleanhtml(raw_html):
  cleanr = re.compile('<.*?>')
  cleantext = re.sub(cleanr, '', raw_html)
  return cleantext
print(cleanhtml(text))


#metadata
doc_result.properties

doc_result.header

doc_result.footer


doc_result.footnotes


doc_html_result = docx2python(path, html = True)

# =============================================================================
# pptx
# =============================================================================


f = open('foobar.pptx')
prs = Presentation(f)
f.close()

# or

with open('foobar.pptx') as f:
    source_stream = StringIO(f.read())
prs = Presentation(source_stream)
source_stream.close()
...
target_stream = StringIO()
prs.save(target_stream)


# =============================================================================
# doc to pdf
# =============================================================================

#linux
from docx2pdf import convert

convert("input.docx")
convert("input.docx", "output.pdf")
convert("my_docx_folder/")
#windows

import sys
import os
import comtypes.client

wdFormatPDF = 17

in_file = os.path.abspath(sys.argv[1])
out_file = os.path.abspath(sys.argv[2])

word = comtypes.client.CreateObject('Word.Application')
doc = word.Documents.Open(in_file)
doc.SaveAs(out_file, FileFormat=wdFormatPDF)
doc.Close()
word.Quit()


# =============================================================================
# 
# =============================================================================
from docx import Document


doc= Document(path) 
for para in doc.paragraphs:
    print(para.text)


tables=[]
for table in doc.tables:
    rows=[]
    for row in table.rows:
        cells=[]
        for cell in row.cells:
            #print(cell.text)
            cells.append(cell.text)
        rows.append(cells)
        df=pd.DataFrame(rows)
        df.columns= df.iloc[0]
        df.drop(0, inplace=True)
        #df.reset_index()
    tables.append(df)


# =============================================================================
# in order
# =============================================================================

from docx import Document
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph

def iter_block_items(parent):
    """
    Generate a reference to each paragraph and table child within *parent*,
    in document order. Each returned value is an instance of either Table or
    Paragraph. *parent* would most commonly be a reference to a main
    Document object, but also works for a _Cell object, which itself can
    contain paragraphs and tables.
    """
    if isinstance(parent, _Document):
        parent_elm = parent.element.body
        # print(parent_elm.xml)
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)

"""
Reading the document.
"""
document = Document(path)
for block in iter_block_items(document):
    #print('found one instance')
    if isinstance(block, Paragraph):
        #print("paragraph")
        print(block.text)
        #write the code here
    else:
        #print("table")
        rows=[]
        for row in block.rows:
            cells=[]
            for cell in row.cells:
                #print(cell.text)
                cells.append(cell.text)
            rows.append(cells)
        print(rows)



def getMetaData(doc):
    metadata = {}
    prop = doc.core_properties
    metadata["author"] = prop.author
    metadata["category"] = prop.category
    metadata["comments"] = prop.comments
    metadata["content_status"] = prop.content_status
    metadata["created"] = prop.created
    metadata["identifier"] = prop.identifier
    metadata["keywords"] = prop.keywords
    metadata["language"] = prop.language
    metadata["modified"] = prop.modified
    metadata["subject"] = prop.subject
    metadata["title"] = prop.title
    metadata["version"] = prop.version
    return metadata

metadata_dict = getMetaData(doc)


# =============================================================================
# ppt
# =============================================================================


path=r"C:\Users\ELECTROBOT\Desktop\da_test_files\VIBE_April 28.pptx"

from pptx import Presentation

ppt= Presentation(path)

text=[]
for num, slide in enumerate(ppt.slides):
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text.append(shape.text)
    
    #dat[num]=text
    
