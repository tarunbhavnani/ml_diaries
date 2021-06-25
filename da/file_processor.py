# -*- coding: utf-8 -*-
"""
Created on Thu Jun 24 13:11:16 2021

@author: ELECTROBOT
"""
import re
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import fitz

from docx import Document

from pptx import Presentation


import glob
files=glob.glob(r"C:\Users\ELECTROBOT\Desktop\da_test_files\*")


file=r"C:\Users\ELECTROBOT\Desktop\da_test_files\rf.pdf"


def tb_index_pdf(file, tb_index):
    doc = fitz.open(file)
    for num, page in enumerate(doc):
        try:
            text = page.getText().encode('utf8')
            text = text.decode('utf8')
            text = clean(text)
            sentences = split_into_sentences(text)
            for sent in sentences:
                tb_index.append({
                    "doc": file.split('\\')[-1],
                    "page": num,
                    "sentence": sent

                })
        except:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": ""

            })
    return tb_index
    
def tb_index_docx(file, tb_index):
    doc= Document(file)
    text=[]
    try:
        for para in doc.paragraphs:
            text.append(para.text)
        text=" ".join([i for i in text if i.strip()!=""])
        #text=text.encode('utf8')
        text = clean(text)
        sentences = split_into_sentences(text)
        for sent in sentences:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": "-",
                "sentence": sent
    
            })
    except:
        tb_index.append({
            "doc": file.split('\\')[-1],
            "page": "-",
            "sentence": "Not read"

        })
    return tb_index



def tb_index_pptx(file, tb_index):
    ppt= Presentation(path)
    
    for num, slide in enumerate(ppt.slides):
        try:
            all_text=[]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text= shape.text
                    text = clean(text)
                    sentences = split_into_sentences(text)
                    [all_text.append(i) for i in sentences]
            all_text= " ".join(i for i in all_text)        
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": all_text
            })

        except:
            tb_index.append({
                "doc": file.split('\\')[-1],
                "page": num,
                "sentence": ""

            })
    return tb_index

                
                
                
                
    


def files_processor_tb(files):
    tb_index = []
    for file in files:
        print(file)
        if file.endswith(".pdf"):
            try:
                tb_index_pdf(file, tb_index)
            except:
                print(file)
        elif file.endswith('.docx'):
            try:
                tb_index_docx(file, tb_index)
            except:
                print(file)
        elif file.endswith('.pptx'):
            try:
                tb_index_pptx(file, tb_index)
            except:
                print(file)
        else:
            print(file)

    
    all_sents = [i['sentence'] for i in tb_index]
    all_sents = all_sents
    vec = TfidfVectorizer(analyzer=ngrams)  # this performs much better but exact words

    vec.fit([i.lower() for i in all_sents])
    
    tfidf_matrix = vec.transform([i.lower() for i in all_sents])
    

    return tb_index, all_sents, vec, tfidf_matrix
